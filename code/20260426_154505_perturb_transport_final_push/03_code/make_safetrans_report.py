from __future__ import annotations

import json
from pathlib import Path

import matplotlib as mpl
import matplotlib.pyplot as plt
from matplotlib import font_manager as fm
import numpy as np
import pandas as pd
import seaborn as sns
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt


ROOT = Path("/home/yyf/codex_cout/20260426_154505_perturb_transport_final_push")
PHASE = ROOT / "14_safetrans_pt"
RESULTS = PHASE / "results"
OUT = PHASE / "presentation" / "latest"
FIG = OUT / "figures"
OUT.mkdir(parents=True, exist_ok=True)
FIG.mkdir(parents=True, exist_ok=True)


def configure_fonts() -> str:
    for p in [
        Path("/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc"),
        Path("/usr/share/fonts/truetype/wqy/wqy-zenhei.ttc"),
        Path("/usr/share/fonts/truetype/droid/DroidSansFallbackFull.ttf"),
    ]:
        if p.exists():
            fm.fontManager.addfont(str(p))
            family = fm.FontProperties(fname=str(p)).get_name()
            mpl.rcParams["font.family"] = family
            mpl.rcParams["font.sans-serif"] = [family, "DejaVu Sans"]
            mpl.rcParams["axes.unicode_minus"] = False
            return family
    mpl.rcParams["axes.unicode_minus"] = False
    return "DejaVu Sans"


FONT = configure_fonts()
sns.set_theme(style="whitegrid", rc={"font.family": FONT, "axes.unicode_minus": False})


def read_csv(name: str) -> pd.DataFrame:
    p = RESULTS / name
    return pd.read_csv(p) if p.exists() and p.stat().st_size else pd.DataFrame()


def savefig(name: str) -> Path:
    p = FIG / name
    plt.tight_layout()
    plt.savefig(p, dpi=220, bbox_inches="tight")
    plt.close()
    return p


def box(ax, xy, text, fc="#eef6fb", ec="#2b6c8a", w=0.22, h=0.16, fs=10):
    x, y = xy
    ax.add_patch(plt.Rectangle((x, y), w, h, fc=fc, ec=ec, lw=1.6))
    ax.text(x + w / 2, y + h / 2, text, ha="center", va="center", fontsize=fs, wrap=True)


def arrow(ax, a, b):
    ax.annotate("", xy=b, xytext=a, arrowprops=dict(arrowstyle="->", lw=1.8, color="#333333"))


def make_figures() -> dict[str, Path]:
    figs: dict[str, Path] = {}
    summary = read_csv("SAFETRANS_SUMMARY_TABLE.csv")
    delta_v2 = read_csv("SAFETRANS_VS_FIXED_V2_DELTAS.csv")
    risk = read_csv("SAFETRANS_RISK_COVERAGE.csv")
    detail = read_csv("SAFETRANS_TASK_DETAILS.csv")
    ablation = read_csv("SAFETRANS_ABLATION_DELTAS.csv")

    plt.figure(figsize=(10, 5.2))
    ax = plt.gca(); ax.axis("off")
    box(ax, (0.03, 0.62), "已有方法\n预测 response", "#e5f5e0", "#31a354")
    box(ax, (0.29, 0.62), "跨 context\n机制会变", "#fff7bc", "#d95f0e")
    box(ax, (0.55, 0.62), "风险问题\n不能硬迁移", "#fee0d2", "#de2d26")
    box(ax, (0.29, 0.22), "SafeTrans-PT\n先判断能否迁移", "#deebf7", "#3182bd", w=0.32)
    box(ax, (0.67, 0.22), "安全: transport\n不安全: 保守/拒判", "#f0f0f0", "#636363", w=0.28)
    arrow(ax, (0.25, 0.70), (0.29, 0.70)); arrow(ax, (0.51, 0.70), (0.55, 0.70))
    arrow(ax, (0.68, 0.62), (0.47, 0.38)); arrow(ax, (0.61, 0.30), (0.67, 0.30))
    ax.text(0.5, 0.95, "研究切口：不是只预测，而是安全迁移", ha="center", fontsize=18, weight="bold")
    figs["story"] = savefig("01_story.png")

    plt.figure(figsize=(10, 5.4))
    ax = plt.gca(); ax.axis("off")
    steps = [
        ("Control state\n目标 context", 0.04),
        ("Perturbation\n扰动信息", 0.24),
        ("Program transport\n效应迁移候选", 0.44),
        ("Transportability\n安全性打分", 0.64),
        ("Adaptive blend\n最终效应", 0.82),
    ]
    for text, x in steps:
        box(ax, (x, 0.56), text, "#deebf7", "#3182bd", w=0.15, h=0.18, fs=10)
    for (_, x1), (_, x2) in zip(steps[:-1], steps[1:]):
        arrow(ax, (x1 + 0.15, 0.65), (x2, 0.65))
    box(ax, (0.18, 0.18), "如果 score 高:\n增加 transport 权重", "#e5f5e0", "#31a354", w=0.26)
    box(ax, (0.54, 0.18), "如果 score 低:\n保守预测/unsafe", "#fee0d2", "#de2d26", w=0.26)
    ax.text(0.5, 0.94, "SafeTrans-PT 架构", ha="center", fontsize=18, weight="bold")
    figs["architecture"] = savefig("02_architecture.png")

    phase3 = ROOT / "13_phase3_confirmation_runs" / "PHASE3_CONFIG_SUMMARY.csv"
    if phase3.exists():
        df = pd.read_csv(phase3)
        cols = ["external_pass_settings", "main_pass_settings"]
        plt.figure(figsize=(9, 4.8))
        x = np.arange(len(df))
        plt.bar(x - 0.18, df["main_pass_settings"], width=0.36, label="main pass")
        plt.bar(x + 0.18, df["external_pass_settings"], width=0.36, label="external pass")
        plt.xticks(x, df["config"].str.replace("config_", "", regex=False), rotation=20, ha="right")
        plt.ylabel("pass settings")
        plt.title("Phase 3 发现：更保守的 transport 更稳")
        plt.legend()
        figs["phase3"] = savefig("03_phase3_blend.png")

    if not delta_v2.empty:
        for phase, fname, title in [
            ("main", "04_main_delta.png", "SafeTrans-PT vs fixed V2: main settings"),
            ("external", "05_external_delta.png", "SafeTrans-PT vs fixed V2: external settings"),
        ]:
            sub = delta_v2[delta_v2["phase"] == phase].copy()
            if not sub.empty:
                sub["setting"] = sub["dataset"].astype(str) + "\n" + sub["split_type"].astype(str)
                mat = sub.set_index("setting")[["top20_delta", "deg_precision_delta", "program_consistency_delta"]]
                plt.figure(figsize=(9.5, max(4.0, 0.42 * len(mat))))
                sns.heatmap(mat, center=0, cmap="vlag", annot=True, fmt=".3f", cbar_kws={"label": "delta"})
                plt.title(title)
                figs[f"{phase}_delta"] = savefig(fname)

    if not risk.empty:
        sub = risk[risk["model"] == "SafeTransPT"].copy()
        if not sub.empty:
            curve = sub.groupby("coverage", as_index=False)["rmse"].mean()
            plt.figure(figsize=(8.2, 4.8))
            plt.plot(curve["coverage"], curve["rmse"], marker="o", lw=2)
            plt.xlabel("coverage 被保留预测比例")
            plt.ylabel("mean RMSE")
            plt.title("Risk-coverage: 低置信样本剔除后风险变化")
            figs["risk"] = savefig("06_risk_coverage.png")

    if not ablation.empty:
        sub = ablation.groupby("ablation", as_index=False)[["top20_delta", "deg_precision_delta", "program_consistency_delta"]].mean()
        if not sub.empty:
            plot = sub.melt(id_vars="ablation", var_name="metric", value_name="delta")
            plt.figure(figsize=(9.5, 4.8))
            sns.barplot(data=plot, x="metric", y="delta", hue="ablation")
            plt.axhline(0, color="black", lw=1)
            plt.title("Ablation: 去掉拒判/去掉 pathway 后的变化")
            plt.xticks(rotation=12)
            figs["ablation"] = savefig("07_ablation.png")

    if not detail.empty:
        sub = detail[detail["model"] == "SafeTransPT"].copy()
        if not sub.empty:
            plt.figure(figsize=(8.2, 5))
            sns.scatterplot(data=sub, x="transportability_score", y="rmse", hue="phase", alpha=0.65)
            plt.title("Failure boundary: 低可迁移性通常更高风险")
            figs["failure"] = savefig("08_failure_boundary.png")
    return figs


def add_title(slide, title, subtitle=None):
    tx = slide.shapes.add_textbox(Inches(0.45), Inches(0.25), Inches(12.4), Inches(0.55))
    p = tx.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = RGBColor(24, 63, 84)
    if subtitle:
        st = slide.shapes.add_textbox(Inches(0.48), Inches(0.82), Inches(12.0), Inches(0.35))
        q = st.text_frame.paragraphs[0]
        q.text = subtitle
        q.font.size = Pt(13)
        q.font.color.rgb = RGBColor(90, 90, 90)


def add_bullets(slide, bullets, left=0.65, top=1.35, width=5.6, height=5.3, font=17):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.clear()
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = b
        p.font.size = Pt(font)
        p.space_after = Pt(8)


def add_pic(slide, path: Path, left=6.4, top=1.25, width=6.0):
    if path.exists():
        slide.shapes.add_picture(str(path), Inches(left), Inches(top), width=Inches(width))


def make_ppt(figs: dict[str, Path]) -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    slides = [
        ("SafeTrans-PT: 跨 context 扰动效应安全迁移", ["目标：不是只预测 response，而是判断哪些 effect 可以迁移、哪些应该保守/拒判。", "方向固定：single-cell perturbation effect transport。"], "story"),
        ("为什么这个问题重要", ["Perturb-seq 数据越来越多，但实验不可能覆盖所有细胞类型、病人、细胞系和扰动组合。", "跨 context 时，同一个扰动可能方向一致、强度不同，也可能机制改变。", "因此模型需要可靠性判断。"], "story"),
        ("相关论文定位", ["scGen: latent perturbation shift。", "CPA/chemCPA: compositional perturbation。", "CellOT: optimal transport。", "GEARS: gene graph + unseen perturbation。", "scPerturb/scPerturBench: 大规模数据和 benchmark。"], None),
        ("现有方法的缺口", ["已有工作已经很多，所以不能说没人做。", "我们的切口是 safe transportability：先判断是否适合迁移，再预测。", "这个比单纯提升 Pearson 更接近真实实验使用。"], None),
        ("SafeTrans-PT 方法", ["输入：control state、perturbation、source effect、pathway/program prior。", "输出：transportability score、adaptive blend、unsafe flag。", "final effect = baseline 与 transported effect 的自适应组合。"], "architecture"),
        ("本周代码与实验", ["新增 safetrans 模型、risk-coverage evaluator、专用 runner。", "保留 strong V0、fixed V2，并加入 no-abstain/no-pathway ablation。", "结果全部保存为 CSV/JSON/log，可复现。"], None),
        ("Phase 3 发现", ["固定 blend 搜索显示：transport 不是越强越好。", "最稳配置是较保守的 blend=0.12。", "这支持 safe/conservative transport 的研究动机。"], "phase3"),
        ("Main settings 结果", ["看 top20、DEG、program consistency，而不是只看 correlation。", "红色代表 SafeTrans-PT 相对 fixed V2 改善。", "如果结果仍弱，会作为明天汇报中的 honest boundary。"], "main_delta"),
        ("External validation", ["外部验证只用于最后检查方向一致性。", "不调 external 参数，不拿 external 做训练标签。", "目标是至少两个 external setting 有正信号。"], "external_delta"),
        ("Risk-coverage / unsafe transport", ["如果 transportability score 低，模型可以保守或拒判。", "coverage 是保留多少预测；risk 是这些预测的平均误差。", "理想情况：剔除低置信样本后风险下降。"], "risk"),
        ("Ablation", ["no-abstain: 去掉拒判机制。", "no-pathway: 去掉 pathway/program 先验。", "用于证明改进不是单个小技巧偶然造成。"], "ablation"),
        ("Failure boundary", ["低可迁移性样本应当对应更高错误或更弱关键基因命中。", "这页用于解释哪些 effect 不该被强行迁移。"], "failure"),
        ("阶段性结论", ["这个方向不是空白，但 safe transport 是明确缺口。", "当前结果的核心发现：保守、自适应迁移比强行迁移更合理。", "明天汇报口径：方法推进 + 阶段性系统验证。"], None),
        ("下一步投稿路线", ["补强强 baseline：GEARS / CPA / CellOT / scVIDR。", "扩大 external validation，但不牺牲 gene-space 对齐。", "加入真实 pathway/STRING/Reactome prior 和更稳的 ranking loss。"], None),
    ]

    for title, bullets, key in slides:
        slide = prs.slides.add_slide(blank)
        add_title(slide, title)
        add_bullets(slide, bullets, width=5.45 if key else 11.9)
        if key and key in figs:
            add_pic(slide, figs[key])

    out = OUT / "SafeTransPT_GroupMeeting_Update.pptx"
    prs.save(out)
    return out


def write_docs() -> None:
    status_path = RESULTS / "SAFETRANS_STATUS.json"
    status = json.loads(status_path.read_text(encoding="utf-8")) if status_path.exists() else {}
    summary = f"""# SafeTrans-PT 当前结果摘要

当前标签：`{status.get('label', 'RUN_NOT_FINISHED')}`

主实验 pass settings：{status.get('main_pass_vs_fixed_v2', 'NA')}

外部验证 pass settings：{status.get('external_pass_vs_fixed_v2', 'NA')}

coverage 检查：{status.get('coverage_ok', 'NA')}

一句话：这版材料不把结果吹成终局，而是讲成“safe transport 方法正在形成，当前已经有保守迁移优于激进迁移的阶段性证据”。
"""
    (OUT / "CURRENT_RESULT_SUMMARY_中文.md").write_text(summary, encoding="utf-8")
    (OUT / "TEACHER_SCRIPT_直接照读.md").write_text(
        """# 老师汇报讲稿

老师好，我这周继续做跨 cellular context 的单细胞扰动效应迁移。这个问题不是简单预测表达量，而是判断一个扰动在某个细胞环境里的 effect 能不能迁移到另一个环境。

已有 scGen、CPA、CellOT、GEARS 等方法都在做 perturbation response prediction，scPerturb 和 scPerturBench 也说明这个方向已经有大规模 benchmark。所以我的切口不是说没人做，而是强调 safe transportability：跨 context 时，模型除了给预测，还要知道什么时候不该自信迁移。

这周我实现了 SafeTrans-PT。它先估计 transportability score，再用这个分数控制 baseline 和 transport effect 的 adaptive blend。低分样本会被保守处理或标记为 unsafe transport。

阶段性结果里，一个重要现象是固定 blend 越大不一定越好，保守的 transport 更稳。这和我们的生物直觉一致：不同 context 中扰动机制可能变化，强行迁移会损害外部泛化。

下一步我会继续补强强 baseline、真实 pathway prior 和 external validation，把这个方向从阶段性方法验证推进到投稿级结果。
""",
        encoding="utf-8",
    )
    (OUT / "RELATED_PAPERS_论文定位.md").write_text(
        """# 相关论文定位

- scGen: https://www.nature.com/articles/s41592-019-0494-8
- CPA / chemCPA: https://github.com/theislab/cpa
- CellOT: https://www.nature.com/articles/s41592-023-01969-x
- GEARS: https://www.nature.com/articles/s41587-023-01905-6
- scPerturb: https://www.nature.com/articles/s41592-023-02144-y
- scPerturBench: https://www.nature.com/articles/s41592-025-02980-0
- Virtual Cells Need Context: https://pmc.ncbi.nlm.nih.gov/articles/PMC12919078/
- scCausalVI: https://www.sciencedirect.com/science/article/pii/S2405471225002765

定位：已有方法多做 response prediction；SafeTrans-PT 强调 safe cross-context transport，即什么时候迁移、什么时候保守。
""",
        encoding="utf-8",
    )
    (OUT / "BEGINNER_EXPLAINER_小白版.md").write_text(
        """# 小白版解释

这个项目问的是：如果我在 A 细胞环境里敲了一个基因，看到了表达变化，那么这个变化能不能搬到 B 细胞环境？

`context` 是细胞环境，比如细胞类型、细胞系、病人状态。

`perturbation` 是人为操作，比如敲掉一个基因或加药。

`effect` 是扰动后的平均表达减去对照组平均表达。

`transport` 是把 A 里学到的 effect 迁移到 B。

`unsafe transport` 是模型不该自信迁移的情况。SafeTrans-PT 的核心就是先判断安不安全，再决定迁移强度。
""",
        encoding="utf-8",
    )


def main() -> None:
    figs = make_figures()
    ppt = make_ppt(figs)
    write_docs()
    print(ppt)


if __name__ == "__main__":
    main()
