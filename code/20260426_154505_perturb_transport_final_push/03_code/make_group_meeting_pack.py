from __future__ import annotations

import shutil
import textwrap
from datetime import datetime
from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path("/home/yyf/codex_cout/20260426_154505_perturb_transport_final_push")
OUT = ROOT / "24_group_meeting_20260514"
FIG_SRC = ROOT / "22_methodology_story_pack" / "figures"
DOCS = OUT / "docs"
FIG = OUT / "figures"
TABLES = OUT / "tables"
PACKAGE = OUT / "package"


def read_csv(path: Path) -> pd.DataFrame:
    return pd.read_csv(path) if path.exists() and path.stat().st_size > 0 else pd.DataFrame()


def mean_by_setting(df: pd.DataFrame) -> pd.DataFrame:
    cols = ["top20_delta", "deg_precision_delta", "program_consistency_delta", "pearson_delta", "spearman_delta"]
    if df.empty:
        return pd.DataFrame()
    rows = []
    for keys, sub in df.groupby(["phase", "split_type"], dropna=False):
        row = dict(zip(["phase", "split_type"], keys))
        row["n"] = len(sub)
        for col in cols:
            row[col] = float(sub[col].mean())
        rows.append(row)
    return pd.DataFrame(rows)


def fmt(v: float) -> str:
    return f"{v:+.3f}"


def add_title(slide, title: str, subtitle: str | None = None):
    box = slide.shapes.add_textbox(Inches(0.45), Inches(0.25), Inches(12.4), Inches(0.7))
    p = box.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = RGBColor(25, 35, 45)
    if subtitle:
        box2 = slide.shapes.add_textbox(Inches(0.48), Inches(0.92), Inches(12.0), Inches(0.35))
        p2 = box2.text_frame.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(12)
        p2.font.color.rgb = RGBColor(90, 100, 110)


def add_bullets(slide, bullets: list[str], x=0.65, y=1.35, w=5.8, h=4.8, size=17):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(size)
        p.space_after = Pt(9)
        p.font.color.rgb = RGBColor(35, 45, 55)


def add_image(slide, path: Path, x, y, w=None, h=None):
    if not path.exists():
        return
    kwargs = {}
    if w is not None:
        kwargs["width"] = Inches(w)
    if h is not None:
        kwargs["height"] = Inches(h)
    slide.shapes.add_picture(str(path), Inches(x), Inches(y), **kwargs)


def add_note(slide, text: str):
    notes = slide.notes_slide.notes_text_frame
    notes.text = textwrap.dedent(text).strip()


def write(path: Path, text: str):
    path.write_text(textwrap.dedent(text).strip() + "\n", encoding="utf-8")


def build_ppt(stats: dict[str, pd.DataFrame]) -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    now = datetime.now().strftime("%Y-%m-%d %H:%M")

    slide = prs.slides.add_slide(blank)
    add_title(slide, "单细胞扰动效应的安全跨环境迁移", "SafeTrans-PT: safe cross-context perturbation effect transport")
    add_bullets(
        slide,
        [
            "核心问题：一个 perturbation effect（扰动效应）能不能从 A 细胞环境迁移到 B 细胞环境？",
            "如果不能迁移，模型应该识别 unsafe transport（不安全迁移），而不是硬预测。",
            f"当前材料更新时间：{now}",
        ],
        x=0.85,
        y=1.75,
        w=11.4,
        size=21,
    )
    add_note(slide, "开场：我这周把问题从普通扰动预测收窄为安全迁移。重点不是说模型万能，而是判断哪些扰动效应可以推广，哪些不能乱推广。")

    slide = prs.slides.add_slide(blank)
    add_title(slide, "为什么这个问题值得做")
    add_bullets(
        slide,
        [
            "单细胞扰动实验很贵，不能每个 cell type（细胞类型）、patient（病人）、cell state（细胞状态）都做一遍。",
            "已有方法多关注 response prediction（扰动响应预测）：预测扰动后表达量像不像。",
            "真实使用时还需要问：这个预测能不能跨 cellular context（细胞环境）信任？",
            "所以我们的目标是：能迁移就预测；不能迁移就标记风险。",
        ],
        x=0.7,
        y=1.25,
        w=6.1,
        size=17,
    )
    add_image(slide, FIG / "01_problem_shift_safe_transport.png", 6.8, 1.35, w=5.8)
    add_note(slide, "这一页强调科学问题：不是换模型刷分，而是解决扰动结论能否推广的问题。")

    slide = prs.slides.add_slide(blank)
    add_title(slide, "已有工作与我们的切口")
    add_bullets(
        slide,
        [
            "scGen（单细胞生成模型）：学习扰动前后的 latent shift（潜在位移）。",
            "CPA（组合扰动自编码器）：建模 drug / dose / context（药物、剂量、环境）的组合响应。",
            "GEARS（基因关系扰动预测）：利用 gene graph（基因图）预测新扰动。",
            "CellOT（最优传输）：学习 control（对照）到 perturbed（扰动后）的分布变化。",
            "我们的切口：transportability（可迁移性），先判断能不能跨 context 安全迁移。",
        ],
        x=0.8,
        y=1.25,
        w=11.6,
        size=16,
    )
    add_note(slide, "这一页不用展开论文细节，只要讲清楚别人主要做预测，我们强调可迁移性和风险边界。")

    slide = prs.slides.add_slide(blank)
    add_title(slide, "数据与验证设置")
    add_bullets(
        slide,
        [
            "使用多个公开单细胞扰动数据集：Haber、Parekh、KaggleCrossCell、KaggleCrossPatient、McFarland、Wessels、Frangieh 等。",
            "补充检查 Norman、Dixit、Papalexi、Tian、Srivatsan 等扰动数据。",
            "held-out perturbation（留出扰动）：训练没见过某些扰动，测试预测它们。",
            "leave-context（留出细胞环境）：训练没见过某个环境，测试迁移到它。",
            "external validation（外部验证）：换独立数据集验证方向是否一致。",
        ],
        x=0.75,
        y=1.18,
        w=11.6,
        size=16,
    )
    add_note(slide, "说明我们不是随机切分，而是用更难、更接近泛化问题的验证方式。")

    slide = prs.slides.add_slide(blank)
    add_title(slide, "方法流程：SafeTrans-PT")
    add_image(slide, FIG / "02_method_pipeline.png", 0.55, 1.12, w=12.1)
    add_note(slide, "照图讲：先整理数据，构造困难划分，然后进入扰动效应空间，再计算可迁移性评分，最后决定预测或保守拒判。")

    slide = prs.slides.add_slide(blank)
    add_title(slide, "当前主结果：DeepSafeTransGPU vs V2")
    main = stats["deep_main"]
    q1 = stats["deep_q1"]
    bullets = []
    if not main.empty:
        hp = main[(main["phase"] == "main") & (main["split_type"] == "heldout_perturbation")]
        ep = main[(main["phase"] == "external") & (main["split_type"] == "heldout_perturbation")]
        if len(hp):
            r = hp.iloc[0]
            bullets.append(f"main held-out perturbation（主数据留出扰动）：top20 {fmt(r.top20_delta)}，DEG {fmt(r.deg_precision_delta)}，program {fmt(r.program_consistency_delta)}")
        if len(ep):
            r = ep.iloc[0]
            bullets.append(f"external held-out perturbation（外部留出扰动）：top20 {fmt(r.top20_delta)}，DEG {fmt(r.deg_precision_delta)}，program {fmt(r.program_consistency_delta)}")
    if not q1.empty:
        ep = q1[(q1["phase"] == "external") & (q1["split_type"] == "heldout_perturbation")]
        if len(ep):
            r = ep.iloc[0]
            bullets.append(f"补强实验外部留出扰动：top20 {fmt(r.top20_delta)}，DEG {fmt(r.deg_precision_delta)}，program {fmt(r.program_consistency_delta)}")
    bullets += [
        "结论：held-out perturbation（留出扰动）是目前最稳的主结果。",
        "注意：这不是说所有跨环境迁移都解决了。",
    ]
    add_bullets(slide, bullets, x=0.75, y=1.18, w=11.7, size=17)
    add_note(slide, "这里是最重要结果页。强调 V2 是更强 baseline，我们已经不是只打赢 V0。")

    slide = prs.slides.add_slide(blank)
    add_title(slide, "结果图：哪些设置更稳")
    add_image(slide, FIG / "03_current_evidence_heatmap.png", 0.65, 1.05, w=12.0)
    add_note(slide, "绿色或正数表示相对 V2 有提升。重点看 held-out perturbation。leave-context 有边界，不要硬吹。")

    slide = prs.slides.add_slide(blank)
    add_title(slide, "社区风格 baseline：补强对照")
    comm = stats["community"]
    bullets = [
        "已补 scGen-style（scGen 风格）、CPA-style（CPA 风格）、OT barycentric（CellOT 风格）、graph kernel（GEARS 风格）轻量对照。",
        "这些不是完整复现原论文，而是用于回答：我们的 baseline 是否太弱？",
    ]
    if not comm.empty:
        ot = comm[(comm.get("model", "") == "OT_barycentric_proxy") & (comm.get("baseline", "") == "V2") & (comm.get("split_type", "") == "heldout_perturbation")]
        if len(ot):
            r = ot.sort_values("phase").iloc[0]
            bullets.append(f"OT 风格 baseline 在 held-out perturbation 也较强：top20 {fmt(r.top20_delta)}，DEG {fmt(r.deg_precision_delta)}。")
    bullets.append("观察：leave-context 下多个 baseline 都明显变差，说明跨环境安全迁移确实是难点。")
    add_bullets(slide, bullets, x=0.75, y=1.18, w=11.7, size=17)
    add_note(slide, "这页的目的不是说我们压倒所有方法，而是说明我们认真补了更强对照，并发现 leave-context 确实难。")

    slide = prs.slides.add_slide(blank)
    add_title(slide, "为什么加入共表达网络模块")
    add_image(slide, FIG / "05_network_module_rationale.png", 0.75, 1.05, w=11.8)
    add_note(slide, "这页讲 hdWGCNA 的借鉴：基因不是单独变，而是一组模块变。模块保守时更可能安全迁移。")

    slide = prs.slides.add_slide(blank)
    add_title(slide, "NetworkSafeTransPT 当前定位")
    net = stats["network"]
    bullets = [
        "NetworkSafeTransPT（网络感知安全迁移）目前不作为主模型。",
        "它对 top20 / DEG 的直接提升不稳定。",
        "但它在 main held-out perturbation 上提升 program/module consistency（基因程序/网络模块一致性）。",
        "所以它更适合做 biological explanation（生物解释层）：解释为什么某些扰动能迁移，某些不能迁移。",
    ]
    if not net.empty:
        hp = net[(net["phase"] == "main") & (net["split_type"] == "heldout_perturbation")]
        if len(hp):
            r = hp.iloc[0]
            bullets.append(f"当前 main held-out perturbation：program {fmt(r.program_consistency_delta_mean)}，two-plus fraction {r.two_plus_effect_fraction:.2f}")
    add_bullets(slide, bullets, x=0.75, y=1.1, w=11.7, size=17)
    add_note(slide, "这一页要诚实：network 不是主结果，而是解释层。这样老师问起来更稳。")

    slide = prs.slides.add_slide(blank)
    add_title(slide, "阶段性结论")
    add_bullets(
        slide,
        [
            "1. 研究问题已收敛：从普通扰动预测，转为扰动效应能否安全跨 context 迁移。",
            "2. 主结果最稳：DeepSafeTransGPU 在 held-out perturbation 上相对 V2 有稳定正向信号。",
            "3. leave-context 仍是难点：不硬说解决，而作为 unsafe transport boundary（不安全迁移边界）。",
            "4. 新增 network module（共表达网络模块）用于解释迁移是否安全。",
            "5. 新增 community-inspired baselines（社区风格基线）用于补强对照。",
        ],
        x=0.75,
        y=1.15,
        w=11.7,
        size=18,
    )
    add_note(slide, "这一页直接总结。")

    slide = prs.slides.add_slide(blank)
    add_title(slide, "下一步计划")
    add_bullets(
        slide,
        [
            "继续收敛主结果：固定 DeepSafeTransGPU 作为主模型，补齐关键数据集和 seed。",
            "把 leave-context 失败案例整理成 unsafe transport case study（不安全迁移案例）。",
            "把 network module preservation（网络模块保守性）做成解释图。",
            "进一步整理社区方法对照，明确我们的优势和边界。",
            "最终写法：安全迁移框架 + 强基线验证 + 生物网络解释。",
        ],
        x=0.75,
        y=1.15,
        w=11.7,
        size=18,
    )
    add_note(slide, "最后一页讲计划，避免把当前结果吹得过满。")

    out = OUT / "SafeTransPT_group_meeting_20260514.pptx"
    prs.save(out)
    return out


def main() -> None:
    for p in [DOCS, FIG, TABLES, PACKAGE]:
        p.mkdir(parents=True, exist_ok=True)
    if FIG_SRC.exists():
        for src in FIG_SRC.glob("*.png"):
            shutil.copy2(src, FIG / src.name)

    deep_main_raw = read_csv(ROOT / "17_may_resume_full_push/reports/ALL_GPU_DEEPSAFE_VS_V2.csv")
    deep_q1_raw = read_csv(ROOT / "20_q1_strong_push/reports/ALL_GPU_DEEPSAFE_VS_V2.csv")
    network = read_csv(ROOT / "21_network_hdWGCNA_push/reports/NETWORK_SAFE_VS_V2_BY_SETTING.csv")
    community = read_csv(ROOT / "23_unique_design_push/reports/COMMUNITY_BASELINE_BY_SETTING.csv")
    deep_main = mean_by_setting(deep_main_raw)
    deep_q1 = mean_by_setting(deep_q1_raw)
    deep_main.to_csv(TABLES / "deep_main_vs_v2_by_setting.csv", index=False)
    deep_q1.to_csv(TABLES / "deep_q1_vs_v2_by_setting.csv", index=False)
    network.to_csv(TABLES / "network_safe_vs_v2_by_setting.csv", index=False)
    community.to_csv(TABLES / "community_baselines_by_setting.csv", index=False)

    stats = {"deep_main": deep_main, "deep_q1": deep_q1, "network": network, "community": community}
    ppt = build_ppt(stats)

    write(
        DOCS / "GROUP_MEETING_README.md",
        f"""
        # 2026-05-14 小组会材料

        ## 最推荐使用

        PPT：`{ppt}`

        ## 汇报主线

        不要讲“我要发几区”。讲科学问题：

        > 单细胞扰动效应能否跨细胞环境安全迁移；如果不能，模型是否能识别不安全迁移。

        ## 当前最稳结论

        DeepSafeTransGPU 在 held-out perturbation（留出扰动）场景下相对 V2 baseline 有稳定正向信号。

        ## 谨慎表述

        leave-context（留出细胞环境）还不能说完全解决，应作为 unsafe transport boundary（不安全迁移边界）。

        ## 材料组成

        - `figures/`：汇报图
        - `tables/`：最新结果表
        - `docs/`：讲稿和说明
        """,
    )
    write(
        DOCS / "SPEAKING_SCRIPT_SHORT_CN.md",
        """
        # 口语讲稿短版

        老师，我这周把问题收敛成了单细胞扰动效应的安全跨环境迁移。

        以前很多扰动预测方法主要关注扰动后表达量能不能预测准，但真实使用时还有一个问题：在一个细胞环境里学到的扰动效应，能不能推广到另一个细胞环境？如果不能，模型应该识别风险，而不是强行预测。

        所以我现在做的 SafeTrans-PT，核心是先判断 transportability，也就是可迁移性，再决定是否进行迁移预测。

        实验上，我用了多个公开单细胞扰动数据集，构造了 held-out perturbation、leave-context 和 external validation。当前最稳定的是 held-out perturbation 场景，DeepSafeTransGPU 相对更强的 V2 baseline 在 top20、DEG precision 和 program consistency 上都有正向结果。

        但 leave-context 仍然更难，所以我不会说已经解决所有跨环境泛化。相反，我把它作为 unsafe transport boundary 来分析，也就是哪些环境下不应该盲目迁移。

        另外，我补了两块内容：一是 community-inspired baselines，用 scGen、CPA、CellOT、GEARS 风格的轻量对照增强可信度；二是借鉴 hdWGCNA 的共表达模块，用 network module preservation 来解释为什么某些扰动效应可以迁移，某些不适合迁移。

        下一步会继续固定主模型，补齐关键数据集和 seed，同时把不安全迁移案例和网络模块解释整理成更完整的图。
        """,
    )

    zip_path = OUT / "SafeTransPT_group_meeting_20260514_pack.zip"
    if zip_path.exists():
        zip_path.unlink()
    if PACKAGE.exists():
        shutil.rmtree(PACKAGE)
    PACKAGE.mkdir(parents=True, exist_ok=True)
    shutil.copy2(ppt, PACKAGE / ppt.name)
    shutil.copytree(DOCS, PACKAGE / "docs")
    shutil.copytree(FIG, PACKAGE / "figures")
    shutil.copytree(TABLES, PACKAGE / "tables")
    shutil.make_archive(str(zip_path.with_suffix("")), "zip", PACKAGE)
    print(zip_path)


if __name__ == "__main__":
    main()
